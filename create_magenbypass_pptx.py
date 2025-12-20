#!/usr/bin/env python3
"""
PowerPoint-Präsentation: Komplikationen des Magenbypasses
Mit medizinischen Diagrammen und Illustrationen
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsmap

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 82, 147)
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(200, 220, 255)
        p2.alignment = PP_ALIGN.CENTER

    return slide

def add_stomach_diagram(slide, left, top, scale=1.0):
    """Add a simplified gastric bypass diagram"""
    # Stomach pouch (small)
    pouch = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top),
        Inches(0.6 * scale), Inches(0.8 * scale)
    )
    pouch.fill.solid()
    pouch.fill.fore_color.rgb = RGBColor(255, 182, 193)  # Light pink
    pouch.line.color.rgb = RGBColor(180, 80, 80)
    pouch.line.width = Pt(2)

    # Bypassed stomach (larger, grayed out)
    bypassed = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.8 * scale), Inches(top + 0.3 * scale),
        Inches(1.0 * scale), Inches(1.4 * scale)
    )
    bypassed.fill.solid()
    bypassed.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Gray
    bypassed.line.color.rgb = RGBColor(150, 150, 150)
    bypassed.line.width = Pt(1)

    # Roux limb (Y-shape intestine)
    roux = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 0.15 * scale), Inches(top + 0.7 * scale),
        Inches(0.3 * scale), Inches(1.5 * scale)
    )
    roux.fill.solid()
    roux.fill.fore_color.rgb = RGBColor(255, 218, 185)  # Peach
    roux.line.color.rgb = RGBColor(180, 120, 80)
    roux.line.width = Pt(1)

    # Biliopancreatic limb
    bp_limb = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 1.0 * scale), Inches(top + 1.6 * scale),
        Inches(0.3 * scale), Inches(1.0 * scale)
    )
    bp_limb.fill.solid()
    bp_limb.fill.fore_color.rgb = RGBColor(255, 218, 185)
    bp_limb.line.color.rgb = RGBColor(180, 120, 80)
    bp_limb.line.width = Pt(1)

    # Common channel
    common = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 0.5 * scale), Inches(top + 2.1 * scale),
        Inches(0.3 * scale), Inches(0.8 * scale)
    )
    common.fill.solid()
    common.fill.fore_color.rgb = RGBColor(255, 218, 185)
    common.line.color.rgb = RGBColor(180, 120, 80)
    common.line.width = Pt(1)

    # Anastomosis marker (red dot)
    anastomosis = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.2 * scale), Inches(top + 0.65 * scale),
        Inches(0.2 * scale), Inches(0.2 * scale)
    )
    anastomosis.fill.solid()
    anastomosis.fill.fore_color.rgb = RGBColor(220, 53, 69)
    anastomosis.line.fill.background()

def add_warning_icon(slide, left, top, size=0.5):
    """Add a warning triangle icon"""
    triangle = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(left), Inches(top),
        Inches(size), Inches(size)
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(255, 193, 7)  # Warning yellow
    triangle.line.color.rgb = RGBColor(200, 150, 0)
    triangle.line.width = Pt(2)

    # Exclamation mark
    exclaim = slide.shapes.add_textbox(
        Inches(left + size * 0.35), Inches(top + size * 0.25),
        Inches(size * 0.3), Inches(size * 0.5)
    )
    tf = exclaim.text_frame
    p = tf.paragraphs[0]
    p.text = "!"
    p.font.size = Pt(int(size * 36))
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

def add_blood_drop_icon(slide, left, top, size=0.4):
    """Add a blood drop icon using oval shape"""
    # Main drop body
    drop = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top + size * 0.3),
        Inches(size), Inches(size)
    )
    drop.fill.solid()
    drop.fill.fore_color.rgb = RGBColor(180, 30, 30)
    drop.line.fill.background()

    # Top point (triangle)
    point = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(left + size * 0.25), Inches(top),
        Inches(size * 0.5), Inches(size * 0.5)
    )
    point.fill.solid()
    point.fill.fore_color.rgb = RGBColor(180, 30, 30)
    point.line.fill.background()

def add_pill_icon(slide, left, top, size=0.5):
    """Add a pill/medication icon"""
    # Pill capsule
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(size * 1.5), Inches(size * 0.6)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(100, 180, 100)
    pill.line.color.rgb = RGBColor(60, 120, 60)
    pill.line.width = Pt(1)

    # Half divider
    half = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + size * 0.75), Inches(top),
        Inches(size * 0.75), Inches(size * 0.6)
    )
    half.fill.solid()
    half.fill.fore_color.rgb = RGBColor(255, 255, 255)
    half.line.fill.background()

def add_heart_icon(slide, left, top, size=0.5):
    """Add a heart icon for cardiovascular"""
    heart = slide.shapes.add_shape(
        MSO_SHAPE.HEART,
        Inches(left), Inches(top),
        Inches(size), Inches(size)
    )
    heart.fill.solid()
    heart.fill.fore_color.rgb = RGBColor(220, 53, 69)
    heart.line.fill.background()

def add_brain_icon(slide, left, top, size=0.6):
    """Add a simplified brain icon (cloud shape)"""
    brain = slide.shapes.add_shape(
        MSO_SHAPE.CLOUD,
        Inches(left), Inches(top),
        Inches(size), Inches(size * 0.8)
    )
    brain.fill.solid()
    brain.fill.fore_color.rgb = RGBColor(255, 182, 193)
    brain.line.color.rgb = RGBColor(180, 100, 120)
    brain.line.width = Pt(1)

def add_checkmark_icon(slide, left, top, size=0.4):
    """Add a checkmark in circle"""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top),
        Inches(size), Inches(size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(40, 167, 69)
    circle.line.fill.background()

    check = slide.shapes.add_textbox(
        Inches(left + size * 0.15), Inches(top),
        Inches(size * 0.7), Inches(size)
    )
    tf = check.text_frame
    p = tf.paragraphs[0]
    p.text = "✓"
    p.font.size = Pt(int(size * 40))
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_intestine_diagram(slide, left, top, scale=1.0):
    """Add intestinal diagram for hernia/stenosis"""
    # Curved intestine segments
    for i in range(4):
        segment = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + (i % 2) * 0.4 * scale),
            Inches(top + i * 0.5 * scale),
            Inches(0.8 * scale), Inches(0.4 * scale)
        )
        segment.fill.solid()
        segment.fill.fore_color.rgb = RGBColor(255, 200, 180)
        segment.line.color.rgb = RGBColor(180, 100, 80)
        segment.line.width = Pt(1)

    # Stenosis marker (narrowing)
    stenosis = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.3 * scale), Inches(top + 0.8 * scale),
        Inches(0.25 * scale), Inches(0.25 * scale)
    )
    stenosis.fill.solid()
    stenosis.fill.fore_color.rgb = RGBColor(220, 53, 69)
    stenosis.line.fill.background()

def add_content_slide(prs, title, bullet_points, icon_type=None):
    """Add a content slide with bullet points and optional icon"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0, 82, 147)
    header.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add icon based on type
    if icon_type == "warning":
        add_warning_icon(slide, 8.8, 1.5, 0.6)
    elif icon_type == "blood":
        add_blood_drop_icon(slide, 8.9, 1.5, 0.5)
        add_blood_drop_icon(slide, 9.15, 1.8, 0.35)
    elif icon_type == "heart":
        add_heart_icon(slide, 8.8, 1.5, 0.7)
    elif icon_type == "pill":
        add_pill_icon(slide, 8.5, 1.6, 0.6)
    elif icon_type == "brain":
        add_brain_icon(slide, 8.6, 1.5, 0.8)
    elif icon_type == "stomach":
        add_stomach_diagram(slide, 7.8, 1.8, 0.7)
    elif icon_type == "intestine":
        add_intestine_diagram(slide, 8.2, 1.5, 0.6)
    elif icon_type == "check":
        add_checkmark_icon(slide, 8.9, 1.5, 0.6)

    # Bullet points
    content_width = Inches(7.8) if icon_type else Inches(8.6)
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), content_width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)
        p.level = 0

    return slide

def add_section_slide(prs, title, number=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 53, 69)
    shape.line.fill.background()

    # Decorative medical cross
    cross_h = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8), Inches(5.5),
        Inches(1.2), Inches(0.4)
    )
    cross_h.fill.solid()
    cross_h.fill.fore_color.rgb = RGBColor(255, 255, 255)
    cross_h.fill.fore_color.brightness = 0.3
    cross_h.line.fill.background()

    cross_v = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.4), Inches(5.1),
        Inches(0.4), Inches(1.2)
    )
    cross_v.fill.solid()
    cross_v.fill.fore_color.rgb = RGBColor(255, 255, 255)
    cross_v.fill.fore_color.brightness = 0.3
    cross_v.line.fill.background()

    # Number
    if number:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Title
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER

    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points, diagram_type=None):
    """Add a two-column content slide with optional diagram"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0, 82, 147)
    header.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Left column box (visual distinction)
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(1.4),
        Inches(4.5), Inches(5.2)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Alice blue
    left_bg.line.color.rgb = RGBColor(0, 82, 147)
    left_bg.line.width = Pt(1)

    # Right column box
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.1), Inches(1.4),
        Inches(4.5), Inches(5.2)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(255, 245, 245)  # Light red
    right_bg.line.color.rgb = RGBColor(220, 53, 69)
    right_bg.line.width = Pt(1)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(4.1), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 82, 147)

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(4.1), Inches(4.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.6), Inches(4.1), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 53, 69)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(2.2), Inches(4.1), Inches(4.2))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)

    return slide

def add_diagram_slide(prs, title, diagram_type, labels):
    """Add a slide with a large diagram and labels"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0, 82, 147)
    header.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if diagram_type == "bypass":
        # Large bypass diagram
        add_detailed_bypass_diagram(slide, 1.5, 1.8, 1.8)

        # Labels on the right
        for i, (label, desc) in enumerate(labels):
            y_pos = 2.0 + i * 0.9
            # Label box
            lbl = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(4), Inches(0.8))
            tf = lbl.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 82, 147)
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(80, 80, 80)

    return slide

def add_detailed_bypass_diagram(slide, left, top, scale):
    """Add a more detailed gastric bypass diagram with labels"""
    # Esophagus
    esoph = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 0.4 * scale), Inches(top - 0.3 * scale),
        Inches(0.25 * scale), Inches(0.5 * scale)
    )
    esoph.fill.solid()
    esoph.fill.fore_color.rgb = RGBColor(255, 200, 180)
    esoph.line.color.rgb = RGBColor(180, 100, 80)

    # Small pouch (active stomach)
    pouch = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.15 * scale), Inches(top + 0.1 * scale),
        Inches(0.75 * scale), Inches(0.9 * scale)
    )
    pouch.fill.solid()
    pouch.fill.fore_color.rgb = RGBColor(255, 150, 150)
    pouch.line.color.rgb = RGBColor(180, 60, 60)
    pouch.line.width = Pt(2)

    # Bypassed stomach (remnant)
    remnant = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 1.1 * scale), Inches(top + 0.2 * scale),
        Inches(1.2 * scale), Inches(1.6 * scale)
    )
    remnant.fill.solid()
    remnant.fill.fore_color.rgb = RGBColor(180, 180, 180)
    remnant.line.color.rgb = RGBColor(120, 120, 120)
    remnant.line.width = Pt(1)

    # Duodenum (bypassed)
    duod = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 1.4 * scale), Inches(top + 1.7 * scale),
        Inches(0.3 * scale), Inches(0.8 * scale)
    )
    duod.fill.solid()
    duod.fill.fore_color.rgb = RGBColor(200, 200, 200)
    duod.line.color.rgb = RGBColor(150, 150, 150)

    # Roux limb (alimentary)
    roux = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 0.3 * scale), Inches(top + 0.9 * scale),
        Inches(0.35 * scale), Inches(1.8 * scale)
    )
    roux.fill.solid()
    roux.fill.fore_color.rgb = RGBColor(255, 220, 200)
    roux.line.color.rgb = RGBColor(200, 140, 100)
    roux.line.width = Pt(2)

    # Biliopancreatic limb
    bp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 1.2 * scale), Inches(top + 2.4 * scale),
        Inches(0.35 * scale), Inches(0.9 * scale)
    )
    bp.fill.solid()
    bp.fill.fore_color.rgb = RGBColor(200, 255, 200)
    bp.line.color.rgb = RGBColor(100, 180, 100)

    # Common channel
    common = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 0.7 * scale), Inches(top + 2.6 * scale),
        Inches(0.35 * scale), Inches(1.0 * scale)
    )
    common.fill.solid()
    common.fill.fore_color.rgb = RGBColor(255, 240, 200)
    common.line.color.rgb = RGBColor(180, 150, 80)

    # Gastrojejunal anastomosis marker
    gj = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.35 * scale), Inches(top + 0.85 * scale),
        Inches(0.25 * scale), Inches(0.25 * scale)
    )
    gj.fill.solid()
    gj.fill.fore_color.rgb = RGBColor(220, 53, 69)
    gj.line.fill.background()

    # JJ anastomosis marker
    jj = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.75 * scale), Inches(top + 2.55 * scale),
        Inches(0.2 * scale), Inches(0.2 * scale)
    )
    jj.fill.solid()
    jj.fill.fore_color.rgb = RGBColor(255, 193, 7)
    jj.line.fill.background()

    # Legend
    legend_y = top + 3.7 * scale

    # Red dot legend
    leg1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(legend_y),
        Inches(0.15), Inches(0.15)
    )
    leg1.fill.solid()
    leg1.fill.fore_color.rgb = RGBColor(220, 53, 69)
    leg1.line.fill.background()

    txt1 = slide.shapes.add_textbox(Inches(left + 0.2), Inches(legend_y - 0.05), Inches(2), Inches(0.3))
    tf = txt1.text_frame
    p = tf.paragraphs[0]
    p.text = "Gastrojejunostomie"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(80, 80, 80)

    # Yellow dot legend
    leg2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(legend_y + 0.25),
        Inches(0.15), Inches(0.15)
    )
    leg2.fill.solid()
    leg2.fill.fore_color.rgb = RGBColor(255, 193, 7)
    leg2.line.fill.background()

    txt2 = slide.shapes.add_textbox(Inches(left + 0.2), Inches(legend_y + 0.2), Inches(2), Inches(0.3))
    tf = txt2.text_frame
    p = tf.paragraphs[0]
    p.text = "Jejunojejunostomie"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(80, 80, 80)

def add_statistics_slide(prs, title, stats):
    """Add a slide with statistics boxes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0, 82, 147)
    header.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Statistics boxes
    colors = [
        RGBColor(0, 82, 147),
        RGBColor(220, 53, 69),
        RGBColor(40, 167, 69),
        RGBColor(255, 193, 7),
        RGBColor(23, 162, 184),
        RGBColor(108, 117, 125)
    ]

    cols = 3
    rows = (len(stats) + cols - 1) // cols

    for i, (value, label) in enumerate(stats):
        row = i // cols
        col = i % cols
        x = 0.5 + col * 3.2
        y = 1.6 + row * 2.4

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(2.9), Inches(2.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Value
        val_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.3), Inches(2.7), Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.1), Inches(2.7), Inches(0.7))
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title with stomach icon
    slide = add_title_slide(prs,
        "Komplikationen des Magenbypasses",
        "Ein umfassender Überblick für medizinisches Fachpersonal")
    # Add decorative medical elements
    add_stomach_diagram(slide, 7.5, 5.2, 0.6)

    # Slide 2: Agenda
    add_content_slide(prs, "Übersicht", [
        "Einführung in den Magenbypass",
        "Klassifikation der Komplikationen",
        "Frühe postoperative Komplikationen",
        "Späte Komplikationen",
        "Ernährungsbedingte Mangelzustände",
        "Psychologische Aspekte",
        "Prävention und Management",
        "Zusammenfassung"
    ], icon_type="check")

    # Slide 3: Anatomy diagram slide
    add_diagram_slide(prs, "Anatomie des Roux-en-Y Magenbypasses", "bypass", [
        ("Magenpouch", "Kleiner verbleibender Magen (~30ml)"),
        ("Restmagen", "Ausgeschalteter Magenanteil (grau)"),
        ("Alimentäre Schlinge", "Roux-Schlinge für Nahrung"),
        ("Biliopankreatische Schlinge", "Führt Verdauungssäfte"),
        ("Common Channel", "Gemeinsame Endstrecke")
    ])

    # Slide 4: Statistics
    add_statistics_slide(prs, "Magenbypass in Zahlen", [
        ("60-80%", "Übergewichtsverlust"),
        ("<0.5%", "Mortalitätsrate"),
        ("10-20%", "Komplikationsrate"),
        ("85%", "Diabetes-Remission"),
        ("70%", "Hypertonie-Besserung"),
        ("5-7 J", "Langzeiterfolg")
    ])

    # Slide 5: Classification
    add_two_column_slide(prs, "Klassifikation der Komplikationen",
        "Frühe Komplikationen (< 30 Tage)",
        ["Anastomoseninsuffizienz", "Blutung", "Lungenembolie",
         "Wundinfektion", "Darmobstruktion", "Atelektase"],
        "Späte Komplikationen (> 30 Tage)",
        ["Dumping-Syndrom", "Nährstoffmängel", "Stenosen",
         "Innere Hernien", "Ulzera", "Cholelithiasis"])

    # Slide 6: Section - Early Complications
    add_section_slide(prs, "Frühe Postoperative Komplikationen", "01")

    # Slide 7: Anastomotic Leak
    add_content_slide(prs, "Anastomoseninsuffizienz", [
        "Inzidenz: 1-5% der Fälle",
        "Häufigste Lokalisation: Gastrojejunostomie",
        "Symptome: Tachykardie (>120/min), Fieber, Bauchschmerzen",
        "Frühzeichen: Unruhe, Tachypnoe, Oligurie",
        "Diagnose: CT mit oralem Kontrastmittel, Methylenblau-Test",
        "Therapie: Drainage, Antibiotika, ggf. Re-Operation",
        "Mortalitätsrisiko: bis zu 15% wenn unbehandelt"
    ], icon_type="warning")

    # Slide 8: Bleeding
    add_content_slide(prs, "Postoperative Blutung", [
        "Inzidenz: 1-4% der Patienten",
        "Intraluminal: GI-Blutung aus Anastomose oder Klammernaht",
        "Extraluminal: Intraabdominelle Blutung (Gefäße, Milz)",
        "Risikofaktoren: Antikoagulation, technische Faktoren",
        "Symptome: Hämatemesis, Meläna, hämodynamische Instabilität",
        "Diagnose: Endoskopie, CT-Angiographie",
        "Therapie: Endoskopische Blutstillung, Transfusion, ggf. Re-OP"
    ], icon_type="blood")

    # Slide 9: Thromboembolism
    add_content_slide(prs, "Thromboembolische Komplikationen", [
        "Tiefe Venenthrombose (TVT): 0,3-1,2%",
        "Lungenembolie (LE): 0,2-1% – Haupttodesursache!",
        "Pathophysiologie: Virchow-Trias verstärkt bei Adipositas",
        "Risikofaktoren: BMI >50, Immobilisation, lange OP-Zeit",
        "Prävention: Frühmobilisation, Kompressionsstrümpfe",
        "Medikamentös: NMH präoperativ beginnen",
        "Therapie: Antikoagulation, ggf. Lyse oder Thrombektomie"
    ], icon_type="heart")

    # Slide 10: Infection
    add_content_slide(prs, "Infektiöse Komplikationen", [
        "Wundinfektion: 2-8% (seltener bei laparoskopischem Zugang)",
        "Intraabdomineller Abszess: 1-2%",
        "Pneumonie: 0,5-2% (Atelektase als Vorstufe)",
        "Risikofaktoren: Diabetes, Immunsuppression, lange OP-Dauer",
        "Prävention: Perioperative Antibiotikaprophylaxe",
        "Wichtig: Gewichtsadaptierte Antibiotikadosierung!",
        "Therapie: Antibiotika, Drainage, Wundpflege"
    ], icon_type="warning")

    # Slide 11: Section - Late Complications
    add_section_slide(prs, "Späte Komplikationen", "02")

    # Slide 12: Dumping Syndrome
    add_two_column_slide(prs, "Dumping-Syndrom",
        "Frühdumping (15-30 min)",
        ["Übelkeit, Erbrechen, Krämpfe",
         "Diarrhoe, Blähungen",
         "Schwitzen, Tachykardie, Flush",
         "Hypotonie, Schwindel",
         "Ursache: Osmotische Flüssigkeitsverschiebung"],
        "Spätdumping (1-3 Stunden)",
        ["Reaktive Hypoglykämie",
         "Schwäche, Zittern, Hunger",
         "Schweißausbrüche, Verwirrtheit",
         "Konzentrationsstörungen",
         "Ursache: Überschießende Insulinsekretion"])

    # Slide 13: Dumping Management
    add_content_slide(prs, "Management des Dumping-Syndroms", [
        "Prävalenz: 20-50% der Bypass-Patienten",
        "Ernährungsmodifikation ist Therapie der ersten Wahl:",
        "   → Kleine, häufige Mahlzeiten (5-6 pro Tag)",
        "   → Reduktion von Einfachzuckern (<15g pro Mahlzeit)",
        "   → Trennung von Essen und Trinken (30 min Abstand)",
        "   → Komplexe Kohlenhydrate und Proteine bevorzugen",
        "Medikamentös: Acarbose 50-100mg zu Mahlzeiten",
        "Therapierefraktär: Octreotid s.c. oder Pasireotid"
    ], icon_type="pill")

    # Slide 14: Section - Nutritional Deficiencies
    add_section_slide(prs, "Ernährungsbedingte Mangelzustände", "03")

    # Slide 15: Vitamin Deficiencies
    add_content_slide(prs, "Vitaminmangelzustände", [
        "Vitamin B12: 30-70% (fehlender Intrinsic-Faktor, Azidität)",
        "   → 1000µg i.m. monatlich oder 1000µg/d oral",
        "Vitamin D: 50-80% (Malabsorption, adipöses Fettgewebe)",
        "   → 3000-6000 IE/d, Ziel: 25-OH-Vit D >30 ng/ml",
        "Folsäure: 15-35% → 400-800µg/d",
        "Vitamin B1 (Thiamin): Cave: Wernicke-Enzephalopathie!",
        "   → Bei Erbrechen: 100mg i.v. vor Glucose!",
        "Vitamin A: 10-15% → Nachtblindheit als Frühzeichen"
    ], icon_type="pill")

    # Slide 16: Mineral Deficiencies
    add_content_slide(prs, "Mineralstoffmängel", [
        "Eisenmangel: 20-50% – häufigste Anämieursache",
        "   → 45-60mg elementares Eisen, getrennt von PPI",
        "Kalziummangel: 10-25% → Osteoporose, Frakturen",
        "   → 1200-1500mg Calciumcitrat + Vitamin D",
        "Zinkmangel: 10-40% → Haarausfall, Wundheilungsstörung",
        "   → 8-22mg/d, bei Mangel 25-50mg/d",
        "Kupfermangel: Selten, aber schwere Myelopathie möglich",
        "Regelmäßige Laborkontrollen alle 3-6 Monate essentiell!"
    ], icon_type="pill")

    # Slide 17: Stenosis
    add_content_slide(prs, "Anastomosenstenose", [
        "Inzidenz: 3-12% an der Gastrojejunostomie",
        "Typische Präsentation: 4-8 Wochen postoperativ",
        "Symptome: Progressive Dysphagie, Übelkeit, Erbrechen",
        "Ursachen: Ischämie, Narbenbildung, Ulkus, Technik",
        "Diagnose: ÖGD mit Strikturnachweis (<10mm Lumen)",
        "Therapie: Endoskopische Ballondilatation (CRE-Ballon)",
        "Prognose: 75% Erfolg nach 1-3 Dilatationen"
    ], icon_type="intestine")

    # Slide 18: Internal Hernia
    add_content_slide(prs, "Innere Hernien", [
        "Inzidenz: 2-5% (häufiger bei laparoskopischem Bypass)",
        "Typen: Petersen-Hernie, Mesenterialschlitz, JJ-Mesentery",
        "Begünstigt durch Gewichtsverlust (weniger Mesenterfett)",
        "Symptome: Intermittierende kolikartige Bauchschmerzen",
        "Gefahr: Darminkarzeration → Ischämie → Nekrose",
        "Diagnose: CT (Wirbelzeichen, geschwollenes Mesenterium)",
        "Therapie: Notfall-OP, meist laparoskopisch möglich"
    ], icon_type="warning")

    # Slide 19: Marginal Ulcer
    add_content_slide(prs, "Marginalulzera (Anastomosenulzera)", [
        "Inzidenz: 1-16% der Patienten",
        "Lokalisation: Jejunale Seite der Gastrojejunostomie",
        "Risikofaktoren: Rauchen (#1!), NSAR, H. pylori",
        "   → Große Pouch, Ischämie, Fremdmaterial",
        "Symptome: Epigastrische Schmerzen, GI-Blutung, Perforation",
        "Diagnose: Ösophagogastroduodenoskopie",
        "Therapie: PPI 2x40mg, Rauchentwöhnung, H. pylori-Eradikation"
    ], icon_type="stomach")

    # Slide 20: Gallstones
    add_content_slide(prs, "Cholelithiasis nach Magenbypass", [
        "Inzidenz: 30-40% entwickeln Gallensteine",
        "Zeitraum: Vor allem in ersten 6-12 Monaten",
        "Ursache: Schneller Gewichtsverlust → Cholesterinsättigung",
        "Symptome: Rechtsseitige Oberbauchschmerzen, Koliken",
        "Prävention: Ursodeoxycholsäure 600mg/d für 6 Monate",
        "Therapie: Laparoskopische Cholezystektomie wenn symptomatisch",
        "Kontrovers: Prophylaktische Cholezystektomie bei OP?"
    ], icon_type="warning")

    # Slide 21: Psychological Complications
    add_content_slide(prs, "Psychologische Komplikationen", [
        "Depression: Erhöhtes Risiko, besonders bei ausbleibendem Erfolg",
        "Suchtverschiebung (Addiction Transfer):",
        "   → Alkohol (beschleunigte Resorption!), Medikamente, Kaufsucht",
        "Essstörungen: Binge Eating, Night Eating Syndrome",
        "Beziehungsprobleme: 20% höhere Scheidungsrate",
        "Suizidalität: 2-4x erhöhtes Risiko (Jahre 2-5 postop)",
        "Prävention: Präoperatives Screening, lebenslange Begleitung"
    ], icon_type="brain")

    # Slide 22: Prevention
    add_content_slide(prs, "Prävention und Langzeitmanagement", [
        "Präoperativ: Sorgfältige Patientenselektion und Edukation",
        "Perioperativ: Thromboseprophylaxe, Antibiotika",
        "Nachsorge: Strukturiertes Programm (lebenslang!)",
        "   → 1., 3., 6. Monat, dann halbjährlich, dann jährlich",
        "Laborkontrollen: BB, Eisen, B12, Folat, Ca, Vit D, PTH",
        "Multivitamin-Supplementierung obligat",
        "Interdisziplinäres Team essentiell"
    ], icon_type="check")

    # Slide 23: Summary Statistics
    add_statistics_slide(prs, "Zusammenfassung: Häufigkeiten", [
        ("1-5%", "Anastomoseninsuffizienz"),
        ("20-50%", "Dumping-Syndrom"),
        ("30-70%", "Vitamin B12-Mangel"),
        ("2-5%", "Innere Hernien"),
        ("1-16%", "Marginalulzera"),
        ("30-40%", "Gallensteine")
    ])

    # Slide 24: Key Takeaways
    add_content_slide(prs, "Kernbotschaften", [
        "Magenbypass ist effektiv, erfordert aber Vigilanz",
        "Frühe Komplikationen: Tachykardie als Warnsignal!",
        "Späte Komplikationen: Oft schleichend, regelmäßige Kontrollen",
        "Nährstoffmängel sind häufig, aber gut behandelbar",
        "Psychische Gesundheit nicht vergessen",
        "Lebenslange Nachsorge ist unerlässlich",
        "Interdisziplinäre Betreuung verbessert Outcomes"
    ], icon_type="check")

    # Slide 25: Thank you
    slide = add_title_slide(prs,
        "Vielen Dank für Ihre Aufmerksamkeit!",
        "Fragen und Diskussion")
    add_heart_icon(slide, 4.7, 5.5, 0.6)

    # Save the presentation
    prs.save('/home/user/Codex-stuff/Komplikationen_Magenbypass.pptx')
    print("Präsentation mit Illustrationen erfolgreich erstellt!")
    print("Datei: Komplikationen_Magenbypass.pptx")
    print(f"Anzahl Folien: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
